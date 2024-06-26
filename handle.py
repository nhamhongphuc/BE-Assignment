import pymupdf
import os
from docx import Document
from docx.shared import Pt, RGBColor
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from pptx import Presentation
from pptx.util import Pt
from googletrans import Translator

### 1 ###
def extract_text_images_from_pdf(pdf_path, output_dir):
    os.makedirs(output_dir, exist_ok=True)
    doc = pymupdf.open(pdf_path)
    all_text = []
    for page in doc:
        text = page.get_text()

        if text.strip():
            all_text.append(text)

        image_list = page.get_images(full=True)
        for img_index, img in enumerate(image_list):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            img_path = os.path.join(output_dir, f"page_{page.number + 1}_img_{img_index + 1}.png")
            with open(img_path, "wb") as img_file:
                img_file.write(image_bytes)

    with open(os.path.join(output_dir, 'extract_text.txt'), 'w', encoding='utf-8') as file:
        for text in all_text:
            file.write(text + '\n')


def extract_text_from_docx(docx_path, output_dir):
    os.makedirs(output_dir, exist_ok=True)
    doc = Document(docx_path)

    all_text = []
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            all_text.append(paragraph.text)
    
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text:
                    all_text.append(cell.text)

    with open(os.path.join(output_dir, 'extract_text.txt'), 'w', encoding='utf-8') as file:
        for text in all_text:
            file.write(text + '\n')

    return all_text

    


### 2 ###
def extract_formatting_from_docx(docx_path, output_dir):
    os.makedirs(output_dir, exist_ok=True)
    doc = Document(docx_path)
    formatting_details = []

    with open(os.path.join(output_dir, 'extract_formatting_details_text.txt'), 'w', encoding='utf-8') as file:
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                para_details = {'text': paragraph.text, 'runs': []}
                file.write(f"Text: {paragraph.text}\n")
                for run in paragraph.runs:
                    font_name = run.font.name if run.font.name else "Default"
                    font_size = run.font.size.pt if run.font.size else "Default"
                    bold = "Yes" if run.bold else "No"
                    italic = "Yes" if run.italic else "No"
                    text_color = run.font.color.rgb if run.font.color and run.font.color.rgb else "Default"
                    file.write(f" - Font: {font_name}, Size: {font_size}, Bold: {bold}, Italic: {italic}, Color: {text_color}\n")
                    run_details = {
                        'text': run.text,
                        'font_name': run.font.name if run.font.name else "Default",
                        'font_size': run.font.size.pt if run.font.size else "Default",
                        'bold': run.bold,
                        'italic': run.italic,
                        'text_color': run.font.color.rgb if run.font.color and run.font.color.rgb else "Default"
                    }
                    para_details['runs'].append(run_details)
                formatting_details.append(para_details)
                file.write("\n")

        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        if paragraph.text.strip():
                            para_details = {'text': paragraph.text, 'runs': []}
                            file.write(f"Text: {paragraph.text}\n")
                            for run in paragraph.runs:
                                font_name = run.font.name if run.font.name else "Default"
                                font_size = run.font.size.pt if run.font.size else "Default"
                                bold = "Yes" if run.bold else "No"
                                italic = "Yes" if run.italic else "No"
                                text_color = run.font.color.rgb if run.font.color and run.font.color.rgb else "Default"
                                file.write(f" - Font: {font_name}, Size: {font_size}, Bold: {bold}, Italic: {italic}, Color: {text_color}\n")
                                run_details = {
                                    'text': run.text,
                                    'font_name': run.font.name if run.font.name else "Default",
                                    'font_size': run.font.size.pt if run.font.size else "Default",
                                    'bold': run.bold,
                                    'italic': run.italic,
                                    'text_color': run.font.color.rgb if run.font.color and run.font.color.rgb else "Default"
                                }
                                para_details['runs'].append(run_details)
                            formatting_details.append(para_details)
                            file.write("\n")
                        file.write("\n")
    return formatting_details

def extract_formatting_from_pdf(pdf_path, output_dir):
    os.makedirs(output_dir, exist_ok=True)
    doc = pymupdf.open(pdf_path)
    
    formatting_details = []
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        blocks = page.get_text("dict")["blocks"]
        
        for block in blocks:
            if block["type"] == 0:
                for line in block["lines"]:
                    for span in line["spans"]:
                        span_details = {
                            'text': span['text'],
                            'font': span['font'],
                            'size': span['size'],
                            'color': span['color'],
                            'page': page_num + 1
                        }
                        formatting_details.append(span_details)
    
    with open(os.path.join(output_dir, 'extract_formatting_details_text.txt'), 'w', encoding='utf-8') as file:
        for detail in formatting_details:
            file.write(f"Page: {detail['page']}\n")
            file.write(f"Text: {detail['text']}\n")
            file.write(f" - Font: {detail['font']}\n")
            file.write(f" - Size: {detail['size']}\n")
            file.write(f" - Color: {detail['color']}\n")
            file.write("\n")
    
    return formatting_details   

def compile_text_from_docx_to_uppercase(docx_path, output_dir):
    text = extract_text_from_docx(docx_path, output_dir)
    formatting = extract_formatting_from_docx(docx_path, output_dir)

    new_doc = Document()
    
    for para in formatting:
        new_para = new_doc.add_paragraph()
        for run in para['runs']:
            new_run = new_para.add_run(run['text'].upper())
            new_run.font.name = run['font_name']
            if run['font_size'] != "Default":
                new_run.font.size = Pt(run['font_size'])
            new_run.bold = run['bold']
            new_run.italic = run['italic']
            if run['text_color'] != "Default":
                new_run.font.color.rgb = RGBColor(run['text_color'][0], run['text_color'][1], run['text_color'][2])
    
    new_doc.save(os.path.join(output_dir, 'output_uppercased_docx_file.docx'))

def compile_text_from_pdf_to_uppercase(pdf_path, output_dir):
    formatting_details = extract_formatting_from_pdf(pdf_path, output_dir)
    output_path = os.path.join(output_dir, 'output_uppercased_pdf.pdf')
    c = canvas.Canvas(output_path, pagesize=letter)
    width, height = letter
    c.setFont("Helvetica", 12)

    y_position = height - 40

    for detail in formatting_details:
        upper_text = detail['text'].upper()
        font_size = detail['size']
        font_name = 'Helvetica'
        color = detail['color']
        
      
        c.setFont(font_name, font_size)
        
      
        r = (color >> 16) & 0xff
        g = (color >> 8) & 0xff
        b = color & 0xff
        c.setFillColorRGB(r/255, g/255, b/255)
        
      
        if y_position < 40:
            c.showPage()
            c.setFont(font_name, font_size)
            y_position = height - 40

        c.drawString(40, y_position, upper_text)
        y_position -= font_size + 4

    c.save()

## 4 ##

def extract_text_images_translate_pptx(pptx_path, output_dir):
    os.makedirs(output_dir, exist_ok=True)
    
    prs = Presentation(pptx_path)
    translator = Translator()

    
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    original_text = paragraph.text
                    try:
                        translated_text = translator.translate(original_text, dest='en').text
                    except Exception as e:
                        translated_text = "Translation Error"
                        print(f"Error translating text: {original_text}, Error: {e}")
                        
                        with open(os.path.join(output_dir, "translation_errors.log"), "a") as error_log:
                            error_log.write(f"Original Text: {original_text}\nError: {e}\n\n")
                    
                    
                    new_run = paragraph.add_run()
                    new_run.text = "\n" + translated_text
                    new_run.font.size = Pt(paragraph.runs[0].font.size.pt if paragraph.runs and paragraph.runs[0].font.size else 12)
                    
                    new_run.font.bold = paragraph.runs[0].font.bold if paragraph.runs else False
                    new_run.font.italic = paragraph.runs[0].font.italic if paragraph.runs else False
                    
            elif shape.shape_type == 13:  
                image = shape.image
                image_bytes = image.blob
                img_path = os.path.join(output_dir, f"slide_{slide_idx + 1}_image_{shape.shape_id}.png")
                with open(img_path, "wb") as img_file:
                    img_file.write(image_bytes)
    
    
    output_pptx_path = os.path.join(output_dir, 'translated_presentation.pptx')
    prs.save(output_pptx_path)
    print(f"Translated presentation saved to {output_pptx_path}")

pdf_path = 'pdf_mock_file.pdf'
pdf_output_dir = './pdf_extract'
doc_path = 'docx_mock_file.docx'
doc_output_dir = './docx_extract'


extract_text_from_docx(doc_path, doc_output_dir)
extract_text_images_from_pdf(pdf_path, pdf_output_dir)

extract_formatting_from_docx(doc_path, doc_output_dir)
extract_formatting_from_pdf(pdf_path, pdf_output_dir)

compile_text_from_docx_to_uppercase(doc_path, doc_output_dir)
compile_text_from_pdf_to_uppercase(pdf_path, pdf_output_dir)

pptx_path = 'Networking.pptx'
output_dir = './pptx_extract'
extract_text_images_translate_pptx(pptx_path, output_dir)